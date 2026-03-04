"""Export Manager Agent - Handles lesson plan and PPT export in multiple formats"""

from typing import Dict, Any, List, Optional
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from docx import Document
from docx.shared import Pt as DocPt
from docx.enum.text import WD_ALIGN_PARAGRAPH
import json
import os
from datetime import datetime
import uuid

class ExportManagerAgent:
    """Agent for exporting lesson plans and PPT presentations in multiple formats"""

    def __init__(self, output_folder: str = "./generated_lessons"):
        """
        Initialize the Export Manager Agent

        Args:
            output_folder: Path to folder where exports will be saved
        """
        self.output_folder = output_folder
        os.makedirs(output_folder, exist_ok=True)

        # Create subfolders for different export types
        self.ppt_folder = os.path.join(output_folder, "pptx")
        self.docx_folder = os.path.join(output_folder, "docx")
        self.json_folder = os.path.join(output_folder, "json")

        os.makedirs(self.ppt_folder, exist_ok=True)
        os.makedirs(self.docx_folder, exist_ok=True)
        os.makedirs(self.json_folder, exist_ok=True)

    async def export_lesson_plan(
        self,
        lesson_data: Dict[str, Any],
        export_format: str = "docx",
        filename_prefix: Optional[str] = None
    ) -> Dict[str, Any]:
        """
        Export lesson plan to specified format

        Args:
            lesson_data: Complete lesson plan data
            export_format: Format to export ('docx', 'json', or 'both')
            filename_prefix: Optional custom filename prefix

        Returns:
            Dictionary containing export results
        """
        try:
            lesson_title = lesson_data.get('lesson_plan', {}).get('title', 'Lesson')
            if not filename_prefix:
                filename_prefix = f"{lesson_title.replace(' ', '_')}_{datetime.now().strftime('%Y%m%d_%H%M%S')}"

            results = {
                "lesson_title": lesson_title,
                "export_timestamp": datetime.now().isoformat(),
                "files": {}
            }

            # Export as Word document
            if export_format in ["docx", "both"]:
                docx_path = await self._export_docx(lesson_data, filename_prefix)
                results["files"]["docx"] = {
                    "path": docx_path,
                    "format": "Word Document (.docx)",
                    "size_mb": os.path.getsize(docx_path) / (1024 * 1024)
                }

            # Export as JSON
            if export_format in ["json", "both"]:
                json_path = self._export_json(lesson_data, filename_prefix)
                results["files"]["json"] = {
                    "path": json_path,
                    "format": "JSON (.json)",
                    "size_mb": os.path.getsize(json_path) / (1024 * 1024)
                }

            return results

        except Exception as e:
            return {
                "error": f"Lesson plan export error: {str(e)}",
                "files": {}
            }

    def _find_best_image_match(
        self,
        query: str,
        indexed_images: List[Dict[str, Any]]
    ) -> Optional[Dict[str, Any]]:
        """
        Find the most semantically relevant image for a query string.
        First tries exact filename match, then falls back to keyword overlap
        against image descriptions.

        Args:
            query: Search string (image description hint or filename)
            indexed_images: List of {path, filename, description} dicts

        Returns:
            Best matching image dict or None
        """
        if not indexed_images:
            return None

        query_lower = query.lower().strip()

        # 1. Exact filename match
        for img in indexed_images:
            if img.get("filename", "").lower() == query_lower:
                return img

        # 2. Filename contains query
        for img in indexed_images:
            if query_lower in img.get("filename", "").lower():
                return img

        # 3. Keyword overlap with description
        query_words = set(query_lower.replace(",", " ").split())
        best_img = None
        best_score = 0
        for img in indexed_images:
            desc_words = set(img.get("description", "").lower().split())
            score = len(query_words & desc_words)
            if score > best_score:
                best_score = score
                best_img = img

        return best_img if best_score > 0 else None

    async def export_ppt(
        self,
        ppt_data: Dict[str, Any],
        template_id: str,
        filename_prefix: Optional[str] = None,
        custom_colors: Optional[Dict[str, str]] = None,
        indexed_images: Optional[List[Dict[str, Any]]] = None
    ) -> Dict[str, Any]:
        """
        Export content as PowerPoint presentation with optional image insertion.

        Args:
            ppt_data: PPT outline and content
            template_id: ID of template to use
            filename_prefix: Optional custom filename prefix
            custom_colors: Optional custom color overrides
            indexed_images: Semantic image index [{path, filename, description}]

        Returns:
            Dictionary containing export results
        """
        try:
            ppt_title = ppt_data.get('ppt_outline', {}).get('title', 'Presentation')
            if not filename_prefix:
                filename_prefix = f"{ppt_title.replace(' ', '_')}_{datetime.now().strftime('%Y%m%d_%H%M%S')}"

            # Create presentation
            prs = Presentation()
            prs.slide_width = Inches(10)
            prs.slide_height = Inches(7.5)

            # Add slides from PPT outline
            slides_data = ppt_data.get('ppt_outline', {}).get('slides', [])

            for slide_info in slides_data:
                slide = self._add_slide_to_presentation(prs, slide_info, custom_colors)
                # --- Semantic image-text alignment ---
                if slide is not None and indexed_images:
                    self._insert_image_into_slide(slide, slide_info, indexed_images)

            # Save presentation
            ppt_path = os.path.join(self.ppt_folder, f"{filename_prefix}.pptx")
            prs.save(ppt_path)

            return {
                "ppt_title": ppt_title,
                "total_slides": len(slides_data),
                "file": {
                    "path": ppt_path,
                    "format": "PowerPoint (.pptx)",
                    "size_mb": os.path.getsize(ppt_path) / (1024 * 1024)
                },
                "export_timestamp": datetime.now().isoformat()
            }

        except Exception as e:
            return {
                "error": f"PPT export error: {str(e)}"
            }

    def _insert_image_into_slide(
        self,
        slide,
        slide_info: Dict[str, Any],
        indexed_images: List[Dict[str, Any]]
    ) -> None:
        """
        Insert the most semantically relevant image into a slide.
        The image is placed in the right half of the slide so it doesn't
        overlap with the text/bullet content on the left.

        Args:
            slide: python-pptx Slide object
            slide_info: Slide content dict (may include 'suggested_images')
            indexed_images: Semantic image index
        """
        from pptx.util import Inches as _Inches
        import os as _os

        # Build a list of candidate queries:
        # 1. LLM-suggested image filenames/descriptions from the slide
        # 2. The slide title as a fallback
        candidate_queries = list(slide_info.get("suggested_images", []))
        if not candidate_queries:
            candidate_queries.append(slide_info.get("title", ""))
        # Add bullet points as additional context for matching
        for bp in slide_info.get("bullet_points", [])[:3]:
            candidate_queries.append(bp)

        matched_image = None
        for query in candidate_queries:
            if not query:
                continue
            matched_image = self._find_best_image_match(query, indexed_images)
            if matched_image:
                break

        if matched_image is None:
            return

        image_path = matched_image.get("path", "")
        if not _os.path.isfile(image_path):
            return

        try:
            # Place image on the right side of the slide
            # Slide dimensions: 10" x 7.5"
            # Right panel: left=5.5", top=1.5", width=4", height=5"
            left = _Inches(5.6)
            top = _Inches(1.4)
            width = _Inches(4.0)
            height = _Inches(5.0)
            slide.shapes.add_picture(image_path, left, top, width=width, height=height)
            print(f"  🖼  Inserted image: {matched_image['filename']} into slide '{slide_info.get('title', '')}'")
        except Exception as e:
            print(f"  ⚠️  Could not insert image {image_path}: {e}")

    async def _export_docx(
        self,
        lesson_data: Dict[str, Any],
        filename_prefix: str
    ) -> str:
        """
        Export lesson plan to Word document

        Args:
            lesson_data: Lesson plan data
            filename_prefix: Filename prefix

        Returns:
            Path to exported file
        """
        doc = Document()

        lesson_plan = lesson_data.get('lesson_plan', {})

        # Add title
        title = doc.add_heading(lesson_plan.get('title', 'Lesson Plan'), 0)
        title.alignment = WD_ALIGN_PARAGRAPH.CENTER

        # Add metadata
        metadata_para = doc.add_paragraph()
        metadata_para.add_run(f"Duration: {lesson_plan.get('duration_minutes', 'N/A')} minutes\n")
        metadata_para.add_run(f"Generated: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}")

        # Add learning objectives
        doc.add_heading('Learning Objectives', level=1)
        for obj in lesson_plan.get('learning_objectives', []):
            doc.add_paragraph(obj, style='List Bullet')

        # Add materials needed
        if lesson_plan.get('materials_needed'):
            doc.add_heading('Materials Needed', level=1)
            for material in lesson_plan.get('materials_needed', []):
                doc.add_paragraph(material, style='List Bullet')

        # Add introduction
        if lesson_plan.get('introduction'):
            doc.add_heading('Introduction', level=1)
            doc.add_paragraph(lesson_plan['introduction'])

        # Add main content
        if lesson_plan.get('main_content'):
            doc.add_heading('Main Content', level=1)
            for section in lesson_plan.get('main_content', []):
                doc.add_heading(section.get('section_title', 'Section'), level=2)
                doc.add_paragraph(section.get('content', ''))

                if section.get('key_points'):
                    doc.add_paragraph('Key Points:')
                    for point in section['key_points']:
                        doc.add_paragraph(point, style='List Bullet')

        # Add activities
        if lesson_plan.get('activities'):
            doc.add_heading('Activities', level=1)
            for activity in lesson_plan.get('activities', []):
                doc.add_heading(activity.get('activity_name', 'Activity'), level=2)
                doc.add_paragraph(activity.get('description', ''))
                doc.add_paragraph(f"Duration: {activity.get('duration_minutes', 'N/A')} minutes")

        # Add assessment
        if lesson_plan.get('assessment'):
            doc.add_heading('Assessment', level=1)
            doc.add_paragraph(lesson_plan['assessment'])

        # Add closure
        if lesson_plan.get('closure'):
            doc.add_heading('Closure', level=1)
            doc.add_paragraph(lesson_plan['closure'])

        # Add discussion questions
        if lesson_data.get('discussion_questions'):
            doc.add_heading('Discussion Questions', level=1)
            for question in lesson_data['discussion_questions']:
                doc.add_paragraph(question, style='List Bullet')

        # Add homework
        if lesson_data.get('homework'):
            doc.add_heading('Homework', level=1)
            doc.add_paragraph(lesson_data['homework'])

        # Save document
        docx_path = os.path.join(self.docx_folder, f"{filename_prefix}.docx")
        doc.save(docx_path)

        return docx_path

    def _export_json(
        self,
        lesson_data: Dict[str, Any],
        filename_prefix: str
    ) -> str:
        """
        Export lesson plan to JSON format

        Args:
            lesson_data: Lesson plan data
            filename_prefix: Filename prefix

        Returns:
            Path to exported file
        """
        # Add export metadata
        export_data = {
            "metadata": {
                "export_timestamp": datetime.now().isoformat(),
                "export_id": str(uuid.uuid4()),
                "version": "1.0"
            },
            "lesson_data": lesson_data
        }

        json_path = os.path.join(self.json_folder, f"{filename_prefix}.json")
        with open(json_path, 'w', encoding='utf-8') as f:
            json.dump(export_data, f, indent=2, ensure_ascii=False)

        return json_path

    def _add_slide_to_presentation(
        self,
        prs: Presentation,
        slide_info: Dict[str, Any],
        custom_colors: Optional[Dict[str, str]] = None
    ):
        """
        Add a slide to the presentation based on slide info.
        Returns the created Slide object so the caller can insert images.
        """
        slide_type = slide_info.get('slide_type', 'content')

        # Create appropriate slide layout
        if slide_type == 'cover':
            slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
            self._add_cover_slide_content(slide, slide_info)
        elif slide_type == 'section':
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            self._add_section_slide_content(slide, slide_info)
        else:  # content, activity, summary
            slide = prs.slides.add_slide(prs.slide_layouts[5])  # Title only layout
            self._add_content_slide(slide, slide_info)

        return slide

    def _add_cover_slide_content(
        self,
        slide,
        slide_info: Dict[str, Any]
    ) -> None:
        """Add content to a cover slide"""
        # Add title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
        title_frame = title_box.text_frame
        title_frame.word_wrap = True
        p = title_frame.paragraphs[0]
        p.text = slide_info.get('title', 'Lesson Plan')
        p.font.size = Pt(54)
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER

        # Add subtitle
        if slide_info.get('subtitle'):
            subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(9), Inches(1))
            subtitle_frame = subtitle_box.text_frame
            subtitle_frame.word_wrap = True
            p = subtitle_frame.paragraphs[0]
            p.text = slide_info['subtitle']
            p.font.size = Pt(32)
            p.alignment = PP_ALIGN.CENTER

    def _add_section_slide_content(
        self,
        slide,
        slide_info: Dict[str, Any]
    ) -> None:
        """Add content to a section slide"""
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(3), Inches(9), Inches(1.5))
        title_frame = title_box.text_frame
        title_frame.word_wrap = True
        p = title_frame.paragraphs[0]
        p.text = slide_info.get('title', 'Section')
        p.font.size = Pt(44)
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER

    def _add_content_slide(
        self,
        slide,
        slide_info: Dict[str, Any]
    ) -> None:
        """Add content to a regular content slide"""
        try:
            # Add title if slide has title placeholder
            if slide.shapes.title:
                slide.shapes.title.text = slide_info.get('title', 'Slide')
        except Exception as e:
            print(f"Warning: Could not set slide title: {e}")

        # Add body text or bullet points
        try:
            # Try to use placeholder 1 (body placeholder)
            if len(slide.placeholders) > 1:
                body_shape = slide.placeholders[1]
                tf = body_shape.text_frame
                tf.clear()
            else:
                # If no body placeholder, create a text box
                body_shape = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(5.5))
                tf = body_shape.text_frame
                tf.word_wrap = True

            if slide_info.get('bullet_points'):
                for i, point in enumerate(slide_info['bullet_points']):
                    if i == 0:
                        if tf.paragraphs:
                            p = tf.paragraphs[0]
                        else:
                            p = tf.add_paragraph()
                    else:
                        p = tf.add_paragraph()
                    p.text = point
                    p.level = 0
                    p.font.size = Pt(18)
            elif slide_info.get('body_text'):
                tf.clear()
                p = tf.add_paragraph()
                p.text = slide_info['body_text']
                p.font.size = Pt(18)
        except Exception as e:
            print(f"Warning: Could not add body content to slide: {e}")
            # Try adding body as a separate text box
            try:
                body_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(5.5))
                body_frame = body_box.text_frame
                body_frame.word_wrap = True

                if slide_info.get('bullet_points'):
                    for i, point in enumerate(slide_info['bullet_points']):
                        p = body_frame.add_paragraph() if i > 0 else body_frame.paragraphs[0]
                        p.text = point
                        p.font.size = Pt(18)
                elif slide_info.get('body_text'):
                    p = body_frame.paragraphs[0]
                    p.text = slide_info['body_text']
                    p.font.size = Pt(18)
            except Exception as e2:
                print(f"Error adding fallback body text: {e2}")

    async def export_full_lesson(
        self,
        lesson_data: Dict[str, Any],
        ppt_data: Dict[str, Any],
        template_id: str,
        lesson_title: str,
        indexed_images: Optional[List[Dict[str, Any]]] = None
    ) -> Dict[str, Any]:
        """
        Export complete lesson (both Word document and PPT)

        Args:
            lesson_data: Lesson plan data
            ppt_data: PPT outline data
            template_id: Template ID to use
            lesson_title: Title for the lesson
            indexed_images: Semantic image index [{path, filename, description}]

        Returns:
            Dictionary containing all export results
        """
        filename_prefix = f"{lesson_title.replace(' ', '_')}_{datetime.now().strftime('%Y%m%d_%H%M%S')}"

        results = {
            "lesson_title": lesson_title,
            "export_id": str(uuid.uuid4()),
            "export_timestamp": datetime.now().isoformat(),
            "exports": {}
        }

        # Export lesson plan
        lesson_export = await self.export_lesson_plan(lesson_data, "both", filename_prefix)
        results["exports"]["lesson_plan"] = lesson_export

        # Export PPT (with image alignment)
        ppt_export = await self.export_ppt(
            ppt_data, template_id, filename_prefix,
            indexed_images=indexed_images
        )
        results["exports"]["ppt"] = ppt_export

        return results


# Example usage for testing
async def test_export_manager():
    """Test the Export Manager Agent"""
    agent = ExportManagerAgent()

    # Sample lesson data
    lesson_data = {
        "lesson_plan": {
            "title": "The Human Circulatory System",
            "learning_objectives": [
                "Understand the structure and function of the heart",
                "Explain how blood circulates through the body"
            ],
            "duration_minutes": 50,
            "materials_needed": ["Diagram of heart", "Video clips"],
            "introduction": "Today we will learn about the human circulatory system...",
            "main_content": [
                {
                    "section_title": "Heart Structure",
                    "content": "The heart is a muscular organ...",
                    "key_points": ["Four chambers", "Pumps blood"]
                }
            ],
            "activities": [],
            "assessment": "Quiz at end of lesson",
            "closure": "Summary and connection to next lesson"
        },
        "discussion_questions": ["Why is the heart important?"],
        "homework": "Read chapter 5"
    }

    # Sample PPT data
    ppt_data = {
        "ppt_outline": {
            "title": "The Human Circulatory System",
            "slides": [
                {
                    "slide_number": 1,
                    "slide_type": "cover",
                    "title": "The Human Circulatory System"
                },
                {
                    "slide_number": 2,
                    "slide_type": "content",
                    "title": "Heart Structure",
                    "bullet_points": ["Left and right chambers", "Pumps blood"]
                }
            ]
        }
    }

    # Export lesson plan
    lesson_export = await agent.export_lesson_plan(lesson_data, "docx", "biology_lesson")
    print("Lesson Export:")
    print(json.dumps(lesson_export, indent=2, ensure_ascii=False))

    # Export PPT
    ppt_export = await agent.export_ppt(ppt_data, "default_formal", "biology_ppt")
    print("\nPPT Export:")
    print(json.dumps(ppt_export, indent=2, ensure_ascii=False))


if __name__ == "__main__":
    import asyncio
    asyncio.run(test_export_manager())
